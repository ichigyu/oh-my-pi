#!/usr/bin/env python3
"""检查 .pptx 模板的可用 layout 和占位符信息，输出结构化 JSON 到 stdout。

用法：
    python3 inspect_template.py template.pptx
"""

import json
import sys


def error_out(category, message, detail=None):
    payload = {
        "ok": False,
        "error": {
            "category": category,
            "message": message,
        },
    }
    if detail is not None:
        payload["error"]["detail"] = detail
    print(json.dumps(payload, ensure_ascii=False))
    sys.exit(1)


try:
    from pptx import Presentation
except ImportError as exc:
    error_out("dependency_error", f"python-pptx 依赖不可用: {exc}")


def inspect(template_path):
    prs = Presentation(template_path)
    layouts = []
    for master in prs.slide_masters:
        for layout in master.slide_layouts:
            placeholders = []
            for ph in layout.placeholders:
                placeholders.append({
                    "idx": ph.placeholder_format.idx,
                    "type": str(ph.placeholder_format.type),
                    "name": ph.name,
                })
            layouts.append({
                "name": layout.name,
                "placeholders": placeholders,
            })

    existing_slides = []
    for i, slide in enumerate(prs.slides):
        slide_placeholders = []
        for ph in slide.placeholders:
            slide_placeholders.append({
                "idx": ph.placeholder_format.idx,
                "type": str(ph.placeholder_format.type),
                "name": ph.name,
            })
        existing_slides.append({
            "index": i,
            "layout_name": slide.slide_layout.name,
            "placeholders": slide_placeholders,
        })

    return {
        "layouts": layouts,
        "existing_slides": existing_slides,
        "slide_size": {"width_emu": prs.slide_width, "height_emu": prs.slide_height},
    }


def main():
    if len(sys.argv) != 2:
        error_out("invalid_args", "usage: inspect_template.py <template.pptx>")

    template_path = sys.argv[1]

    try:
        result = inspect(template_path)
    except Exception as exc:
        error_out("inspect_error", str(exc))

    print(json.dumps({"ok": True, **result}, ensure_ascii=False, indent=2))


if __name__ == "__main__":
    main()
