#!/usr/bin/env python3
"""从 JSON 输入生成 .pptx 文件。

JSON 输入来自 stdin，或通过 --input <path> 传入文件。
支持模板填充模式（传入 template）和无模板从零生成模式。
输出结构化 JSON 到 stdout：{"ok": true, "output": "..."} 或 {"ok": false, "error": {...}}。
"""

import argparse
import json
import os
import sys
import traceback


def error_out(category, message, detail=None):
    payload = {
        "ok": False,
        "error": {
            "category": category,
            "message": message,
        },
    }
    if detail is not None:
        payload["error"]["detail"] = detail
    print(json.dumps(payload, ensure_ascii=False))
    sys.exit(1)


try:
    from pptx import Presentation
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE
    from pptx.util import Emu, Inches, Pt
except ImportError as exc:
    error_out("dependency_error", f"python-pptx 依赖不可用: {exc}")


CHART_TYPES = {
    "bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "line": XL_CHART_TYPE.LINE,
    "pie": XL_CHART_TYPE.PIE,
}


def find_layout(prs, layout_name):
    for master in prs.slide_masters:
        for layout in master.slide_layouts:
            if layout.name == layout_name:
                return layout
    raise ValueError(f"layout 不存在: {layout_name}（可用 inspect_template.py 查看可用 layout）")


def find_placeholder(slide, key):
    """按占位符名称或 idx 匹配；key 可以是名称字符串或数字字符串（idx）。"""
    for ph in slide.placeholders:
        if ph.name == key:
            return ph
    if key.isdigit():
        idx = int(key)
        for ph in slide.placeholders:
            if ph.placeholder_format.idx == idx:
                return ph
    return None


def fill_text_placeholder(ph, value):
    if ph.has_text_frame:
        ph.text_frame.text = str(value)


def fill_table_placeholder(slide, ph, table_data):
    rows = len(table_data)
    cols = len(table_data[0]) if rows else 0
    left, top, width, height = ph.left, ph.top, ph.width, ph.height
    graphic_frame = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = graphic_frame.table
    for r, row in enumerate(table_data):
        for c, cell_value in enumerate(row):
            table.cell(r, c).text = str(cell_value)


def fill_picture_placeholder(slide, ph, image_path):
    if not os.path.isfile(image_path):
        raise ValueError(f"图片文件不存在: {image_path}")
    left, top, width, height = ph.left, ph.top, ph.width, ph.height
    slide.shapes.add_picture(image_path, left, top, width=width, height=height)


def add_chart_to_slide(slide, chart_spec):
    chart_type = CHART_TYPES.get(chart_spec.get("type"))
    if chart_type is None:
        raise ValueError(f"unsupported chart type: {chart_spec.get('type')} (支持 {list(CHART_TYPES)})")

    categories = chart_spec.get("categories", [])
    series = chart_spec.get("series", [])
    if not categories or not series:
        raise ValueError("chart spec missing required fields: categories, series")

    chart_data = CategoryChartData()
    chart_data.categories = categories
    for s in series:
        chart_data.add_series(s.get("name", "Series"), s.get("values", []))

    left = Emu(chart_spec.get("left_emu", Inches(1)))
    top = Emu(chart_spec.get("top_emu", Inches(1.5)))
    width = Emu(chart_spec.get("width_emu", Inches(8)))
    height = Emu(chart_spec.get("height_emu", Inches(4.5)))

    slide.shapes.add_chart(chart_type, left, top, width, height, chart_data)


def fill_slide(slide, slide_spec):
    placeholders_spec = slide_spec.get("placeholders", {})
    for key, value in placeholders_spec.items():
        ph = find_placeholder(slide, key)
        if ph is None:
            raise ValueError(f"占位符不存在: {key}")

        if isinstance(value, dict):
            if "image" in value:
                fill_picture_placeholder(slide, ph, value["image"])
            elif "table" in value:
                fill_table_placeholder(slide, ph, value["table"])
            elif "text" in value:
                fill_text_placeholder(ph, value["text"])
            else:
                raise ValueError(f"占位符 {key} 的值缺少 image/table/text 字段")
        else:
            fill_text_placeholder(ph, value)

    for chart_spec in slide_spec.get("charts", []):
        add_chart_to_slide(slide, chart_spec)


def generate_from_template(spec):
    template_path = spec.get("template")
    output_path = spec.get("output")
    slides_spec = spec.get("slides")

    if not os.path.isfile(template_path):
        raise ValueError(f"template 文件不存在: {template_path}")
    if not isinstance(slides_spec, list) or not slides_spec:
        raise ValueError("missing required field: slides (must be non-empty list)")

    prs = Presentation(template_path)

    for slide_spec in slides_spec:
        if not isinstance(slide_spec, dict):
            raise ValueError("invalid slide spec: 每个 slide 必须是 JSON object")
        layout_name = slide_spec.get("layout")
        if not layout_name:
            raise ValueError("invalid slide spec: 缺少 layout 字段")
        layout = find_layout(prs, layout_name)
        slide = prs.slides.add_slide(layout)
        fill_slide(slide, slide_spec)

    prs.save(output_path)
    return output_path


def generate_without_template(spec):
    output_path = spec.get("output")
    slides_spec = spec.get("slides")

    if not isinstance(slides_spec, list) or not slides_spec:
        raise ValueError("missing required field: slides (must be non-empty list)")

    prs = Presentation()

    for slide_spec in slides_spec:
        if not isinstance(slide_spec, dict):
            raise ValueError("invalid slide spec: 每个 slide 必须是 JSON object")
        layout_idx = slide_spec.get("layout_index", 1)
        if not isinstance(layout_idx, int) or layout_idx < 0 or layout_idx >= len(prs.slide_layouts):
            raise ValueError(f"invalid layout_index: {layout_idx}")
        layout = prs.slide_layouts[layout_idx]
        slide = prs.slides.add_slide(layout)
        fill_slide(slide, slide_spec)

    prs.save(output_path)
    return output_path


def generate(spec):
    if not isinstance(spec, dict):
        raise ValueError("invalid spec: 根节点必须是 JSON object")

    output_path = spec.get("output")
    if not output_path:
        raise ValueError("missing required field: output")

    if spec.get("template"):
        output_path = generate_from_template(spec)
    else:
        output_path = generate_without_template(spec)

    if not os.path.isfile(output_path):
        raise RuntimeError(f"保存后文件不存在: {output_path}")
    try:
        Presentation(output_path)
    except Exception as exc:
        raise RuntimeError(f"保存后 python-pptx 无法读取生成的文件: {exc}")

    return output_path


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--input", help="JSON 输入文件路径；不传则从 stdin 读取")
    args = parser.parse_args()

    try:
        if args.input:
            with open(args.input, "r", encoding="utf-8") as f:
                raw = f.read()
        else:
            raw = sys.stdin.read()
    except OSError as exc:
        error_out("io_error", f"无法读取输入: {exc}")

    try:
        spec = json.loads(raw)
    except json.JSONDecodeError as exc:
        error_out("invalid_json", f"JSON 解析失败: {exc}")

    try:
        output_path = generate(spec)
    except ValueError as exc:
        error_out("invalid_spec", str(exc))
    except Exception as exc:
        error_out("generation_error", str(exc), traceback.format_exc())

    print(json.dumps({"ok": True, "output": output_path}, ensure_ascii=False))


if __name__ == "__main__":
    main()
